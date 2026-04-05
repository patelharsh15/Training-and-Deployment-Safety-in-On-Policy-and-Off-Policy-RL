#!/usr/bin/env python3
"""Generate a proper .pptx PowerPoint presentation with all images embedded."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

PLOTS = os.path.join(os.path.dirname(os.path.dirname(__file__)), "plots")
OUT = os.path.join(os.path.dirname(__file__), "presentation.pptx")

# Colors
DARK_BLUE = RGBColor(0x1a, 0x23, 0x7e)
MED_BLUE = RGBColor(0x28, 0x35, 0x93)
LIGHT_BG = RGBColor(0xe8, 0xea, 0xf6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
RED = RGBColor(0xd3, 0x2f, 0x2f)
GREEN = RGBColor(0x2e, 0x7d, 0x32)
ORANGE_BG = RGBColor(0xFF, 0xF3, 0xE0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_paragraph(tf, text, font_size=18, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, space_before=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    p.alignment = alignment
    p.space_before = space_before
    return p

def add_image(slide, img_name, left, top, width=None, height=None):
    path = os.path.join(PLOTS, img_name)
    if not os.path.exists(path):
        print(f"WARNING: {path} not found, skipping.")
        return
    if width and height:
        slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width), Inches(height))
    elif width:
        slide.shapes.add_picture(path, Inches(left), Inches(top), width=Inches(width))
    elif height:
        slide.shapes.add_picture(path, Inches(left), Inches(top), height=Inches(height))
    else:
        slide.shapes.add_picture(path, Inches(left), Inches(top))

def add_rect(slide, left, top, width, height, fill_color=LIGHT_BG, border_color=DARK_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    return shape

###############################################################################
# SLIDE 1: TITLE
###############################################################################
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, WHITE)
# Title bar
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
shape.fill.solid(); shape.fill.fore_color.rgb = DARK_BLUE; shape.line.fill.background()
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.35), prs.slide_width, Inches(0.15))
shape.fill.solid(); shape.fill.fore_color.rgb = DARK_BLUE; shape.line.fill.background()

add_text_box(slide, 1, 1.8, 11.3, 1.5, "Training and Deployment Safety in\nOn-Policy and Off-Policy RL",
             font_size=40, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.8, 11.3, 0.6, "Farina Salman  •  Rachna Sunilkumar Deshpande  •  Abdulaziz Al-Tayar  •  Harsh Patel",
             font_size=20, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 4.5, 11.3, 0.5, "University of Ottawa — CSI 5180 / CEG 5270: Topics in AI",
             font_size=18, color=RGBColor(0x88, 0x88, 0x88), alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 5.1, 11.3, 0.5, "April 2026",
             font_size=18, color=RGBColor(0x88, 0x88, 0x88), alignment=PP_ALIGN.CENTER)

###############################################################################
# SLIDE 2: AGENDA
###############################################################################
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11.7, 0.7, "Agenda", font_size=36, bold=True, color=DARK_BLUE)

agenda_items = [
    "1.  Motivation & Problem Statement",
    "2.  Gap in the Literature",
    "3.  Research Questions",
    "4.  Experimental Methodology",
    "5.  Results — MuJoCo Environments (Hopper-v4, HalfCheetah-v4)",
    "6.  Results — Safety-Gymnasium (SafetyPointGoal1-v0)",
    "7.  Qualitative Video Analysis",
    "8.  Discussion & Algorithmic Insights",
    "9.  Limitations & Future Work",
    "10. Conclusion & References",
]
tf = add_text_box(slide, 2, 1.5, 9, 5, agenda_items[0], font_size=22, color=DARK_GRAY)
for item in agenda_items[1:]:
    add_paragraph(tf, item, font_size=22, color=DARK_GRAY, space_before=Pt(12))

###############################################################################
# SLIDE 3: MOTIVATION
###############################################################################
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11.7, 0.7, "1. Motivation", font_size=36, bold=True, color=DARK_BLUE)

bullets = [
    "• RL is being deployed in safety-critical domains: autonomous vehicles, robotic surgery, industrial control",
    "• PPO (on-policy) and SAC (off-policy) achieve high rewards, but behave very differently in terms of stability, safety, and reliability",
    "• An unsafe action in a real-world system can cause physical damage, injury, or catastrophic failure",
    "• We need to understand which algorithm is safer before deploying RL in the real world",
]
tf = add_text_box(slide, 0.8, 1.5, 7, 3.5, bullets[0], font_size=20, color=DARK_GRAY)
for b in bullets[1:]:
    add_paragraph(tf, b, font_size=20, color=DARK_GRAY, space_before=Pt(14))

rect = add_rect(slide, 8.5, 1.5, 4.2, 2.5)
tf2 = rect.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = "Core Question:"
p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = DARK_BLUE
add_paragraph(tf2, "If we give PPO and SAC the exact same compute budget, environments, and evaluation protocol — which one produces policies that are safer to deploy on a physical robot?", font_size=16, color=DARK_GRAY, space_before=Pt(10))

###############################################################################
# SLIDE 4: GAP IN LITERATURE
###############################################################################
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11.7, 0.7, "2. Gap in the Literature", font_size=36, bold=True, color=DARK_BLUE)

add_text_box(slide, 0.8, 1.3, 5.5, 0.4, "What prior work covers:", font_size=20, bold=True, color=MED_BLUE)
prior = [
    "• Cumulative reward comparisons [Schulman'17, Haarnoja'18]",
    "• Sample efficiency benchmarks [Henderson'18]",
    "• Implementation details [Engstrom'20]",
    "• On-policy best practices [Andrychowicz'21]",
]
tf = add_text_box(slide, 0.8, 1.8, 5.5, 2.5, prior[0], font_size=18, color=DARK_GRAY)
for p in prior[1:]:
    add_paragraph(tf, p, font_size=18, color=DARK_GRAY, space_before=Pt(8))

add_text_box(slide, 7, 1.3, 5.5, 0.4, "What is missing:", font_size=20, bold=True, color=RED)
missing = [
    "• Training and deployment safety comparison",
    "• Policy smoothness under matched conditions",
    "• Constraint adherence in safety-critical envs",
    "• Recovery from perturbations",
    "• Risk sensitivity analysis",
]
tf = add_text_box(slide, 7, 1.8, 5.5, 2.5, missing[0], font_size=18, color=RED)
for m in missing[1:]:
    add_paragraph(tf, m, font_size=18, color=RED, space_before=Pt(8))

rect = add_rect(slide, 0.8, 4.8, 11.7, 1.5)
tf3 = rect.text_frame; tf3.word_wrap = True
p = tf3.paragraphs[0]
p.text = "Our Contribution: A direct, reproducible safety comparison of PPO vs SAC using practical deployment metrics (action smoothness, constraint violations, recovery, gradient stability) across both standard and constrained environments."
p.font.size = Pt(18); p.font.color.rgb = DARK_BLUE; p.font.bold = False

###############################################################################
# SLIDE 5: RESEARCH QUESTIONS
###############################################################################
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11.7, 0.7, "3. Research Questions", font_size=36, bold=True, color=DARK_BLUE)

rect = add_rect(slide, 0.8, 1.5, 11.7, 1.5)
tf = rect.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Primary RQ:"
p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = DARK_BLUE
add_paragraph(tf, "In which environmental conditions will PPO have more stable and safer behavior compared to SAC for both training and deployment, when they both have the same compute budgets and evaluation protocols?", font_size=18, color=DARK_GRAY, space_before=Pt(8))

rect2 = add_rect(slide, 0.8, 3.5, 11.7, 1.5)
tf2 = rect2.text_frame; tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = "Secondary RQ:"
p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = DARK_BLUE
add_paragraph(tf2, "What are the algorithmic mechanisms that cause the observed differences in stability and safety between PPO and SAC?", font_size=18, color=DARK_GRAY, space_before=Pt(8))

add_text_box(slide, 0.8, 5.5, 11.7, 1, "We evaluate the secondary RQ through gradient norm analysis, Q-value variance tracking, and policy entropy evolution — linking behavioral differences to on-policy vs. off-policy mechanics.",
             font_size=16, color=RGBColor(0x88, 0x88, 0x88))

###############################################################################
# SLIDE 6: METHODOLOGY
###############################################################################
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11.7, 0.7, "4. Experimental Methodology", font_size=36, bold=True, color=DARK_BLUE)

rows = [
    ("Component", "Details"),
    ("Algorithms", "PPO (on-policy) vs SAC (off-policy) via Stable-Baselines3"),
    ("Network", "MLP Policy & Value: 2×256 ReLU"),
    ("Environments", "Hopper-v4, HalfCheetah-v4 (MuJoCo) + SafetyPointGoal1-v0 (Safety-Gym)"),
    ("Baseline", "Random policy (per environment)"),
    ("Budget", "1,000,000 environment steps per run (matched)"),
    ("Seeds", "10 independent seeds per condition (60 total runs)"),
    ("Parallelization", "8× SubprocVecEnv (vectorized environments)"),
    ("Hardware", "2× NVIDIA RTX 3070 (8GB VRAM), 36 CPU cores"),
    ("HP Search", "LR ∈ {1e-4, 3e-4, 1e-3}, γ ∈ {0.99, 0.995}, BS ∈ {256, 1024}"),
    ("Evaluation", "100 deterministic episodes per seed"),
]
from pptx.util import Inches as In
tbl = slide.shapes.add_table(len(rows), 2, In(0.8), In(1.3), In(11.7), In(5.5)).table
tbl.columns[0].width = In(2.5)
tbl.columns[1].width = In(9.2)
for i, (c1, c2) in enumerate(rows):
    for j, val in enumerate([c1, c2]):
        cell = tbl.cell(i, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(16)
            p.font.name = 'Calibri'
            if i == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_BLUE
            else:
                p.font.color.rgb = DARK_GRAY
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_BG

###############################################################################
# SLIDE 7: SAFETY METRICS
###############################################################################
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11.7, 0.7, "Safety & Stability Metrics", font_size=36, bold=True, color=DARK_BLUE)

metrics_rows = [
    ("Metric", "What It Measures", "Why It Matters"),
    ("Action Smoothness", "L2 norm of consecutive action differences", "Jerky actions = mechanical wear on real robots"),
    ("Constraint Violations", "Hazard boundary breaches per episode", "Direct measure of safety compliance"),
    ("Feasibility Rate", "% episodes with zero violations", "Reliability of safe operation"),
    ("Recovery Ratio", "Performance after state perturbation", "Robustness to unexpected disturbances"),
    ("Policy Entropy", "Randomness of action selection", "Exploration vs. exploitation balance"),
    ("Q-Value Variance", "Uncertainty of value estimates (SAC)", "Confidence in decision making"),
    ("Gradient Norms", "Magnitude of parameter updates", "Training stability indicator"),
]
tbl = slide.shapes.add_table(len(metrics_rows), 3, In(0.5), In(1.3), In(12.3), In(5.5)).table
tbl.columns[0].width = In(2.5)
tbl.columns[1].width = In(5)
tbl.columns[2].width = In(4.8)
for i, row in enumerate(metrics_rows):
    for j, val in enumerate(row):
        cell = tbl.cell(i, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(15)
            p.font.name = 'Calibri'
            if i == 0:
                p.font.bold = True; p.font.color.rgb = WHITE
                cell.fill.solid(); cell.fill.fore_color.rgb = DARK_BLUE
            else:
                p.font.color.rgb = DARK_GRAY
                if i % 2 == 0:
                    cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_BG

###############################################################################
# SLIDE 8: SUMMARY TABLE
###############################################################################
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11.7, 0.7, "5. Results Overview", font_size=36, bold=True, color=DARK_BLUE)
add_image(slide, "summary_table.png", 1.5, 1.3, width=10)
add_text_box(slide, 0.8, 6.2, 11.7, 0.5, "Mean ± standard error across 10 seeds, evaluated over 100 deterministic episodes each.",
             font_size=14, color=RGBColor(0x88, 0x88, 0x88), alignment=PP_ALIGN.CENTER)

###############################################################################
# SLIDE 9: HOPPER LEARNING CURVES
###############################################################################
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11.7, 0.7, "Hopper-v4 — Learning Curves", font_size=36, bold=True, color=DARK_BLUE)
add_image(slide, "learning_curves_hopper_v4.png", 1, 1.2, width=7.5)
rect = add_rect(slide, 8.8, 1.5, 4, 2.5)
tf = rect.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Key Finding:"
p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = DARK_BLUE
add_paragraph(tf, "PPO (~2744) and SAC (~2762) achieve nearly identical final returns. Both algorithms solve this balance task equally well within the 1M step budget.", font_size=15, color=DARK_GRAY, space_before=Pt(10))

###############################################################################
# SLIDE 10: HOPPER SMOOTHNESS
###############################################################################
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11.7, 0.7, "Hopper-v4 — Action Smoothness", font_size=36, bold=True, color=DARK_BLUE)
add_image(slide, "action_smoothness_hopper_v4.png", 0.5, 1.2, width=7)

add_text_box(slide, 8, 1.5, 2.2, 0.4, "PPO Smoothness", font_size=14, color=RGBColor(0x66, 0x66, 0x66), alignment=PP_ALIGN.CENTER)
add_text_box(slide, 8, 1.9, 2.2, 0.6, "0.129", font_size=36, bold=True, color=GREEN, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 10.5, 1.5, 2.2, 0.4, "SAC Smoothness", font_size=14, color=RGBColor(0x66, 0x66, 0x66), alignment=PP_ALIGN.CENTER)
add_text_box(slide, 10.5, 1.9, 2.2, 0.6, "0.397", font_size=36, bold=True, color=RED, alignment=PP_ALIGN.CENTER)

rect = add_rect(slide, 8, 3, 4.8, 2.5)
tf = rect.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "PPO's actions are 3× smoother than SAC's. On a physical robot, SAC's jerky motor commands would cause severe mechanical wear. PPO's trust-region clipping naturally produces stable, deployment-safe control signals."
p.font.size = Pt(16); p.font.color.rgb = DARK_GRAY

###############################################################################
# SLIDE 11: HOPPER GRADIENT NORMS
###############################################################################
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11.7, 0.7, "Hopper-v4 — Gradient Norm Evolution", font_size=36, bold=True, color=DARK_BLUE)
add_image(slide, "gradient_norms_envsteps_hopper_v4.png", 0.8, 1.2, width=8)
rect = add_rect(slide, 0.8, 5.8, 11.7, 1.2)
tf = rect.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Algorithmic Insight: PPO's gradient norms are flat and stable (clipped trust region). SAC's gradients spike violently (up to 2500), reflecting its aggressive entropy-maximizing exploration. This directly explains the 3× smoothness gap."
p.font.size = Pt(16); p.font.color.rgb = DARK_BLUE

###############################################################################
# SLIDE 12: HOPPER SEED VARIANCE
###############################################################################
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11.7, 0.7, "Hopper-v4 — Seed Variance", font_size=36, bold=True, color=DARK_BLUE)
add_image(slide, "seed_variance_hopper_v4.png", 0.5, 1.2, width=7)
rect = add_rect(slide, 8, 1.5, 4.8, 3)
tf = rect.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Both algorithms show consistent performance across all 10 seeds, confirming reproducibility."
p.font.size = Pt(16); p.font.color.rgb = DARK_GRAY
add_paragraph(tf, "SAC has a slightly tighter IQR (±254 vs ±326), suggesting more deterministic convergence on this simple task.", font_size=16, color=DARK_GRAY, space_before=Pt(14))

###############################################################################
# SLIDE 13: HALFCHEETAH LEARNING CURVES
###############################################################################
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11.7, 0.7, "HalfCheetah-v4 — Learning Curves", font_size=36, bold=True, color=DARK_BLUE)
add_image(slide, "learning_curves_halfcheetah_v4.png", 1, 1.2, width=7.5)
rect = add_rect(slide, 8.8, 1.5, 4, 2.5)
tf = rect.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Key Finding:"
p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = DARK_BLUE
add_paragraph(tf, "SAC (~6270) dominates PPO (~1138) by 5.5×. In an environment where failure is impossible, SAC's aggressive exploration pays off massively.", font_size=15, color=DARK_GRAY, space_before=Pt(10))

###############################################################################
# SLIDE 14: REWARD HACKING
###############################################################################
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11.7, 0.7, "HalfCheetah-v4 — Reward Hacking Phenomenon", font_size=36, bold=True, color=DARK_BLUE)
add_image(slide, "action_smoothness_halfcheetah_v4.png", 0.5, 1.2, width=7)

rect = add_rect(slide, 8, 1.3, 4.8, 2.2, fill_color=ORANGE_BG, border_color=RGBColor(0xe6, 0x51, 0x00))
tf = rect.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "⚠ Reward Hacking Detected!"
p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = RGBColor(0xe6, 0x51, 0x00)
add_paragraph(tf, "PPO's agent learned to run completely upside down, sliding on its head. Its trust region prevented risky corrections, trapping it in this local optimum (~1100).", font_size=15, color=DARK_GRAY, space_before=Pt(10))

rect2 = add_rect(slide, 8, 3.8, 4.8, 2.2)
tf2 = rect2.text_frame; tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = "SAC's entropy maximization forced it to keep trying chaotic movements, eventually discovering the upright running gait (~6270). This is a textbook stability-exploration tradeoff."
p.font.size = Pt(15); p.font.color.rgb = DARK_GRAY

###############################################################################
# SLIDE 15: HALFCHEETAH GRADIENT NORMS
###############################################################################
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11.7, 0.7, "HalfCheetah-v4 — Gradient Norms & Q-Value Variance", font_size=36, bold=True, color=DARK_BLUE)
add_image(slide, "gradient_norms_envsteps_halfcheetah_v4.png", 0.8, 1.2, width=8)
rect = add_rect(slide, 0.8, 5.8, 11.7, 1.2)
tf = rect.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "SAC's Q-value variance on HalfCheetah: 3.44 (vs 0.27 on Hopper — a 13× increase). The critic network is far less certain, yet this very uncertainty drives better exploration and dramatically higher performance."
p.font.size = Pt(16); p.font.color.rgb = DARK_BLUE

###############################################################################
# SLIDE 16: HALFCHEETAH SEED VARIANCE
###############################################################################
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11.7, 0.7, "HalfCheetah-v4 — Seed Variance", font_size=36, bold=True, color=DARK_BLUE)
add_image(slide, "seed_variance_halfcheetah_v4.png", 0.5, 1.2, width=7)
rect = add_rect(slide, 8, 1.5, 4.8, 3.5)
tf = rect.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "PPO shows remarkably low variance (±55) because all 10 seeds converged to the same local optimum (upside-down sliding)."
p.font.size = Pt(16); p.font.color.rgb = DARK_GRAY
add_paragraph(tf, "SAC has higher variance (±176) reflecting different exploration paths, but consistently achieves the superior upright gait.", font_size=16, color=DARK_GRAY, space_before=Pt(14))

###############################################################################
# SLIDE 17: SAFETY LEARNING CURVES
###############################################################################
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11.7, 0.7, "6. SafetyPointGoal1-v0 — Learning Curves", font_size=36, bold=True, color=DARK_BLUE)
add_image(slide, "learning_curves_safetypointgoal1_v0.png", 1, 1.2, width=7.5)
rect = add_rect(slide, 8.8, 1.5, 4, 2.5)
tf = rect.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Both PPO and SAC achieve identical returns (~27.0). The real differentiation is in constraint violations, not reward."
p.font.size = Pt(16); p.font.color.rgb = DARK_GRAY

###############################################################################
# SLIDE 18: SAFETY CONSTRAINTS
###############################################################################
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11.7, 0.7, "SafetyPointGoal1-v0 — Constraint Violations", font_size=36, bold=True, color=DARK_BLUE)
add_image(slide, "constraints_safetypointgoal1_v0.png", 0.3, 1.2, width=7)

add_text_box(slide, 8, 1.5, 2.2, 0.4, "PPO Mean Cost", font_size=14, color=RGBColor(0x66, 0x66, 0x66), alignment=PP_ALIGN.CENTER)
add_text_box(slide, 8, 1.9, 2.2, 0.6, "54.6", font_size=36, bold=True, color=RED, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 10.5, 1.5, 2.2, 0.4, "SAC Mean Cost", font_size=14, color=RGBColor(0x66, 0x66, 0x66), alignment=PP_ALIGN.CENTER)
add_text_box(slide, 10.5, 1.9, 2.2, 0.6, "50.3", font_size=36, bold=True, color=GREEN, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 8, 2.8, 2.2, 0.4, "PPO Feasibility", font_size=14, color=RGBColor(0x66, 0x66, 0x66), alignment=PP_ALIGN.CENTER)
add_text_box(slide, 8, 3.2, 2.2, 0.5, "8%", font_size=30, bold=True, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 10.5, 2.8, 2.2, 0.4, "SAC Feasibility", font_size=14, color=RGBColor(0x66, 0x66, 0x66), alignment=PP_ALIGN.CENTER)
add_text_box(slide, 10.5, 3.2, 2.2, 0.5, "9%", font_size=30, bold=True, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

rect = add_rect(slide, 8, 4.2, 4.8, 2)
tf = rect.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Surprise Finding: SAC achieves ~8% fewer violations than PPO. SAC's replay buffer allows it to memorize hazard locations from past experiences."
p.font.size = Pt(16); p.font.color.rgb = DARK_BLUE

###############################################################################
# SLIDE 19: REWARD-SAFETY TRADEOFF
###############################################################################
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11.7, 0.7, "SafetyPointGoal1-v0 — Reward-Safety Tradeoff", font_size=36, bold=True, color=DARK_BLUE)
add_image(slide, "reward_safety_safetypointgoal1_v0.png", 1, 1.2, width=7.5)
rect = add_rect(slide, 8.8, 1.5, 4, 3)
tf = rect.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Neither algorithm successfully learned to avoid hazards while maximizing reward. Both cluster in the high-cost region, highlighting the fundamental difficulty of safety-constrained RL without explicit cost penalties."
p.font.size = Pt(16); p.font.color.rgb = DARK_GRAY

###############################################################################
# SLIDE 20: SAFETY SMOOTHNESS
###############################################################################
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11.7, 0.7, "SafetyPointGoal1-v0 — Action Smoothness", font_size=36, bold=True, color=DARK_BLUE)
add_image(slide, "action_smoothness_safetypointgoal1_v0.png", 0.5, 1.2, width=7)

add_text_box(slide, 8, 1.5, 2.2, 0.4, "PPO Smoothness", font_size=14, color=RGBColor(0x66, 0x66, 0x66), alignment=PP_ALIGN.CENTER)
add_text_box(slide, 8, 1.9, 2.2, 0.6, "0.256", font_size=36, bold=True, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 10.5, 1.5, 2.2, 0.4, "SAC Smoothness", font_size=14, color=RGBColor(0x66, 0x66, 0x66), alignment=PP_ALIGN.CENTER)
add_text_box(slide, 10.5, 1.9, 2.2, 0.6, "0.135", font_size=36, bold=True, color=GREEN, alignment=PP_ALIGN.CENTER)

rect = add_rect(slide, 8, 3, 4.8, 2.5)
tf = rect.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Interestingly, SAC produces smoother actions on the Safety env. The navigation task rewards precise movements toward the goal — and SAC's replay buffer helps it refine these smooth trajectories over time."
p.font.size = Pt(16); p.font.color.rgb = DARK_GRAY

###############################################################################
# SLIDE 21: SAFETY GRADIENT NORMS
###############################################################################
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11.7, 0.7, "SafetyPointGoal1-v0 — Gradient Norms", font_size=36, bold=True, color=DARK_BLUE)
add_image(slide, "gradient_norms_envsteps_safetypointgoal1_v0.png", 0.8, 1.2, width=8)
rect = add_rect(slide, 0.8, 5.8, 11.7, 1.2)
tf = rect.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Both algorithms show relatively stable gradient norms on this simpler navigation task, but SAC's occasional spikes correlate with episodes where the agent encounters novel hazard configurations."
p.font.size = Pt(16); p.font.color.rgb = DARK_BLUE

###############################################################################
# SLIDE 22: RECOVERY
###############################################################################
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11.7, 0.7, "Recovery from Perturbations", font_size=36, bold=True, color=DARK_BLUE)
add_image(slide, "recovery_comparison.png", 0.5, 1.2, width=7)

rect = add_rect(slide, 8, 1.5, 4.8, 2.5)
tf = rect.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Recovery Ratio > 1.0 means the agent performed better after perturbation."
p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = DARK_BLUE
add_paragraph(tf, "All conditions show ratios above 1.0, indicating both algorithms are robust to unexpected state disturbances — a positive sign for deployment.", font_size=15, color=DARK_GRAY, space_before=Pt(10))

# Recovery table
rows = [("Environment", "PPO", "SAC"), ("Hopper-v4", "1.96", "1.94"), ("HalfCheetah-v4", "2.34", "2.30"), ("SafetyPointGoal1", "1.83", "1.67")]
tbl = slide.shapes.add_table(len(rows), 3, In(8), In(4.5), In(4.8), In(2)).table
for i, row in enumerate(rows):
    for j, val in enumerate(row):
        cell = tbl.cell(i, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14); p.font.name = 'Calibri'
            if i == 0:
                p.font.bold = True; p.font.color.rgb = WHITE
                cell.fill.solid(); cell.fill.fore_color.rgb = DARK_BLUE
            else:
                p.font.color.rgb = DARK_GRAY

###############################################################################
# SLIDE 23: SAMPLE EFFICIENCY
###############################################################################
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11.7, 0.7, "Sample Efficiency", font_size=36, bold=True, color=DARK_BLUE)
add_image(slide, "sample_efficiency.png", 2, 1.2, width=9)
rect = add_rect(slide, 0.8, 5.8, 11.7, 1.2)
tf = rect.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "SAC demonstrates superior sample efficiency on HalfCheetah, reaching performance thresholds significantly faster due to its replay buffer enabling data reuse. On Hopper, both algorithms converge at similar rates."
p.font.size = Pt(16); p.font.color.rgb = DARK_BLUE

###############################################################################
# SLIDE 24: VIDEO ANALYSIS
###############################################################################
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11.7, 0.7, "7. Qualitative Video Analysis", font_size=36, bold=True, color=DARK_BLUE)

vid_rows = [
    ("Video", "Observed Behavior", "Metric Correlation"),
    ("hopper-v4_ppo", "Stable, repetitive, cautious hopping rhythm", "Smoothness: 0.129 ✓"),
    ("hopper-v4_sac", "Wild, explosive micro-adjustments", "Smoothness: 0.397 ✗"),
    ("cheetah_ppo", "Running completely upside down!", "Return: 1138 (local opt.)"),
    ("cheetah_sac", "Upright, fast sprinting gait", "Return: 6270 (global opt.)"),
    ("safety_ppo", "Clips hazard boundaries frequently", "Cost: 54.6 violations"),
    ("safety_sac", "Course-corrects faster near hazards", "Cost: 50.3 violations"),
]
tbl = slide.shapes.add_table(len(vid_rows), 3, In(0.5), In(1.3), In(12.3), In(5)).table
tbl.columns[0].width = In(2.5)
tbl.columns[1].width = In(5.5)
tbl.columns[2].width = In(4.3)
for i, row in enumerate(vid_rows):
    for j, val in enumerate(row):
        cell = tbl.cell(i, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(16); p.font.name = 'Calibri'
            if i == 0:
                p.font.bold = True; p.font.color.rgb = WHITE
                cell.fill.solid(); cell.fill.fore_color.rgb = DARK_BLUE
            else:
                p.font.color.rgb = DARK_GRAY
                if i % 2 == 0:
                    cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_BG

add_text_box(slide, 0.8, 6.5, 11.7, 0.5, "All videos rendered at 30 fps using headless EGL GPU rendering on RTX 3070.",
             font_size=13, color=RGBColor(0x88, 0x88, 0x88))

###############################################################################
# SLIDE 25: DISCUSSION
###############################################################################
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11.7, 0.7, "8. Discussion — Answering the Research Questions", font_size=34, bold=True, color=DARK_BLUE)

rect = add_rect(slide, 0.8, 1.3, 11.7, 2.2)
tf = rect.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Primary RQ — When is PPO safer?"
p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = DARK_BLUE
add_paragraph(tf, "PPO is safer in physics-based balance tasks (Hopper) where its trust-region clipping produces 3× smoother actions. However, this same conservatism traps PPO in local optima (upside-down HalfCheetah), and in constrained navigation, SAC actually achieves fewer violations.", font_size=16, color=DARK_GRAY, space_before=Pt(10))

rect2 = add_rect(slide, 0.8, 3.8, 11.7, 3)
tf2 = rect2.text_frame; tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = "Secondary RQ — What mechanisms cause the differences?"
p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = DARK_BLUE
add_paragraph(tf2, "• PPO's clipped surrogate objective → stable gradients → smooth actions → but risk of local optima entrapment", font_size=16, color=DARK_GRAY, space_before=Pt(12))
add_paragraph(tf2, "• SAC's entropy maximization → volatile gradients → jerky actions → but escapes local optima", font_size=16, color=DARK_GRAY, space_before=Pt(8))
add_paragraph(tf2, "• SAC's replay buffer → data reuse → faster learning of spatial constraints → fewer safety violations", font_size=16, color=DARK_GRAY, space_before=Pt(8))

###############################################################################
# SLIDE 26: KEY TAKEAWAY
###############################################################################
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11.7, 0.7, "The Core Insight", font_size=36, bold=True, color=DARK_BLUE)

rect = add_rect(slide, 1.5, 1.3, 10.3, 2, fill_color=ORANGE_BG, border_color=RGBColor(0xe6, 0x51, 0x00))
tf = rect.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "There is no universally \"safer\" algorithm."
p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = RGBColor(0xe6, 0x51, 0x00); p.alignment = PP_ALIGN.CENTER
add_paragraph(tf, "PPO excels at deployment safety (smooth, predictable actions)", font_size=18, color=DARK_GRAY, alignment=PP_ALIGN.CENTER, space_before=Pt(14))
add_paragraph(tf, "SAC excels at constraint compliance (fewer boundary violations)", font_size=18, color=DARK_GRAY, alignment=PP_ALIGN.CENTER, space_before=Pt(6))
add_paragraph(tf, "The choice depends on what kind of safety matters most.", font_size=18, bold=True, color=DARK_GRAY, alignment=PP_ALIGN.CENTER, space_before=Pt(6))

winner_rows = [
    ("Safety Dimension", "Winner", "Evidence"),
    ("Action Smoothness (Hopper)", "PPO ✓", "3× smoother (0.129 vs 0.397)"),
    ("Exploration Robustness", "SAC ✓", "Avoids local optima (6270 vs 1138)"),
    ("Constraint Violations", "SAC ✓", "8% fewer violations (50.3 vs 54.6)"),
    ("Gradient Stability", "PPO ✓", "Flat norms vs SAC's violent spikes"),
    ("Recovery from Perturbation", "Tie", "Both > 1.0 across all environments"),
]
tbl = slide.shapes.add_table(len(winner_rows), 3, In(1), In(3.8), In(11.3), In(3)).table
tbl.columns[0].width = In(4)
tbl.columns[1].width = In(2)
tbl.columns[2].width = In(5.3)
for i, row in enumerate(winner_rows):
    for j, val in enumerate(row):
        cell = tbl.cell(i, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(15); p.font.name = 'Calibri'
            if i == 0:
                p.font.bold = True; p.font.color.rgb = WHITE
                cell.fill.solid(); cell.fill.fore_color.rgb = DARK_BLUE
            else:
                p.font.color.rgb = DARK_GRAY
                if j == 1 and "PPO" in val:
                    p.font.color.rgb = GREEN; p.font.bold = True
                elif j == 1 and "SAC" in val:
                    p.font.color.rgb = MED_BLUE; p.font.bold = True

###############################################################################
# SLIDE 27: LIMITATIONS
###############################################################################
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11.7, 0.7, "9. Limitations & Future Work", font_size=36, bold=True, color=DARK_BLUE)

add_text_box(slide, 0.8, 1.3, 5.5, 0.4, "Limitations", font_size=22, bold=True, color=RED)
lim = [
    "• Only one safety environment tested",
    "• Standard SB3 implementations — no constrained RL (CPO, FOCOPS)",
    "• Fixed 1M step budget — longer training may change conclusions",
    "• MLP architecture only — no CNN/attention",
    "• No real-world hardware validation (sim-only)",
]
tf = add_text_box(slide, 0.8, 1.9, 5.5, 4, lim[0], font_size=18, color=DARK_GRAY)
for l in lim[1:]:
    add_paragraph(tf, l, font_size=18, color=DARK_GRAY, space_before=Pt(10))

add_text_box(slide, 7, 1.3, 5.5, 0.4, "Future Work", font_size=22, bold=True, color=GREEN)
fw = [
    "• Add constrained RL algorithms (CPO, FOCOPS, SafePPO)",
    "• Test more Safety-Gym environments (Car, Doggo)",
    "• Increase budget to 5–10M steps",
    "• Sim-to-real transfer on physical robot",
    "• Lagrangian reward shaping for cost minimization",
    "• Investigate reward hacking mitigation",
]
tf = add_text_box(slide, 7, 1.9, 5.5, 4, fw[0], font_size=18, color=DARK_GRAY)
for f in fw[1:]:
    add_paragraph(tf, f, font_size=18, color=DARK_GRAY, space_before=Pt(10))

###############################################################################
# SLIDE 28: CONCLUSION
###############################################################################
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11.7, 0.7, "10. Conclusion", font_size=36, bold=True, color=DARK_BLUE)

conclusions = [
    "• Performed a rigorous, reproducible safety comparison of PPO vs SAC across 3 environments, 10 seeds each, with 7 custom safety metrics",
    "• PPO produces 3× smoother actions — critical for robotic deployment",
    "• SAC achieves 5.5× higher performance in unconstrained environments",
    "• SAC shows fewer constraint violations in safety-critical environments thanks to its replay buffer",
    "• Discovered a reward hacking phenomenon in PPO (HalfCheetah) — a direct consequence of over-conservative trust-region optimization",
    "• Recommendation: Use PPO when actuator wear matters; use SAC when constraint compliance matters",
]
tf = add_text_box(slide, 0.8, 1.3, 11.7, 4, conclusions[0], font_size=20, color=DARK_GRAY)
for c in conclusions[1:]:
    add_paragraph(tf, c, font_size=20, color=DARK_GRAY, space_before=Pt(14))

rect = add_rect(slide, 2, 5.8, 9.3, 1)
tf = rect.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Code & Data: github.com/patelharsh15/Training-and-Deployment-Safety-in-On-Policy-and-Off-Policy-RL"
p.font.size = Pt(16); p.font.color.rgb = DARK_BLUE; p.alignment = PP_ALIGN.CENTER

###############################################################################
# SLIDE 29: REFERENCES
###############################################################################
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11.7, 0.7, "References", font_size=36, bold=True, color=DARK_BLUE)

refs = [
    '[1] Schulman et al., "Proximal Policy Optimization Algorithms," arXiv:1707.06347, 2017.',
    '[2] Haarnoja et al., "Soft Actor-Critic: Off-Policy Maximum Entropy Deep RL," ICML, 2018.',
    '[3] Henderson et al., "Deep Reinforcement Learning That Matters," AAAI, 2018.',
    '[4] Engstrom et al., "Implementation Matters in Deep RL: A Case Study on PPO and TRPO," ICLR, 2020.',
    '[5] Andrychowicz et al., "What Matters in On-Policy RL? A Large-Scale Empirical Study," ICLR, 2021.',
    '[6] Ray et al., "Benchmarking Safe Exploration in Deep RL," OpenAI Technical Report, 2019.',
    '[7] Ji et al., "Safety-Gymnasium: A Unified Safe RL Benchmark," NeurIPS, 2023.',
    '[8] Achiam et al., "Constrained Policy Optimization," ICML, 2017.',
]
tf = add_text_box(slide, 0.8, 1.3, 11.7, 5, refs[0], font_size=16, color=DARK_GRAY)
for r in refs[1:]:
    add_paragraph(tf, r, font_size=16, color=DARK_GRAY, space_before=Pt(10))

###############################################################################
# SLIDE 30: THANK YOU
###############################################################################
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
shape.fill.solid(); shape.fill.fore_color.rgb = DARK_BLUE; shape.line.fill.background()
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.35), prs.slide_width, Inches(0.15))
shape.fill.solid(); shape.fill.fore_color.rgb = DARK_BLUE; shape.line.fill.background()

add_text_box(slide, 1, 2, 11.3, 1, "Thank You!", font_size=48, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.2, 11.3, 0.8, "Questions?", font_size=36, color=MED_BLUE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 4.5, 11.3, 0.5, "Farina Salman  •  Rachna Sunilkumar Deshpande  •  Abdulaziz Al-Tayar  •  Harsh Patel",
             font_size=20, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 5.2, 11.3, 0.5, "University of Ottawa — CSI 5180 / CEG 5270",
             font_size=18, color=RGBColor(0x88, 0x88, 0x88), alignment=PP_ALIGN.CENTER)

rect = add_rect(slide, 3, 6, 7.3, 0.8)
tf = rect.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "GitHub: github.com/patelharsh15/Training-and-Deployment-Safety-in-On-Policy-and-Off-Policy-RL"
p.font.size = Pt(15); p.font.color.rgb = DARK_BLUE; p.alignment = PP_ALIGN.CENTER

###############################################################################
# SAVE
###############################################################################
prs.save(OUT)
print(f"\n✅ Presentation saved to: {OUT}")
print(f"   Total slides: {len(prs.slides)}")
